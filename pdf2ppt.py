import tkinter as tk
from tkinter import filedialog, messagebox, Entry, Button
import webbrowser
from ttkbootstrap import Style
from PIL import ImageTk, Image
import PyPDF2
from pptx import Presentation
from transformers import pipeline

class PDFSummarizerGUI:
    def __init__(self, master):
        self.master = master
        self.master.title("PDF2PPT Summarizer")
        self.master.iconbitmap("Assets\logo.ico")

        # Set up ttkbootstrap style with solar theme
        self.style = Style(theme='solar')
        
        # Load background image
        img = Image.open("Assets\gui.png")
        img = img.resize((self.master.winfo_screenwidth(), self.master.winfo_screenheight()), resample=Image.LANCZOS)
        bg_img = ImageTk.PhotoImage(img)
        
        # Set background image as label
        bg_label = tk.Label(self.master, image=bg_img)
        bg_label.place(x=0, y=0, relwidth=1, relheight=1)
        bg_label.image = bg_img

        main_frame = tk.Frame(self.master)
        main_frame.pack()

        self.label1 = tk.Label(main_frame, text="Select a PDF file to summarize:")
        self.label1.pack(pady=10)

        self.filepath_entry = Entry(main_frame, width=50)
        self.filepath_entry.pack(pady=10)

        self.browse_button = Button(main_frame, text="Browse", command=self.browse_pdf)
        self.browse_button.pack(pady=10)

        self.summarize_button = Button(main_frame, text="Summarize", command=self.summarize_pdf)
        self.summarize_button.pack(pady=10)
        
        self.attribution_label = tk.Label(self.master, text="Made by Michael Assefa", font=("Times New Roman", 10))
        self.attribution_label.pack(side="bottom", pady=10)
        
        self.linked_in_label = tk.Label(main_frame, text="LinkedIn", fg="blue", cursor="hand2", font=("Times New Roman", 10)) 
        self.linked_in_label.pack(side="bottom", padx=10, pady=10)
        self.linked_in_label.bind("<Button-1>", lambda e: webbrowser.open_new("https://www.linkedin.com/in/michael-assefa-965643221/"))
        
        self.github_label = tk.Label(main_frame, text="GitHub", fg="blue", cursor="hand2", font=("Times New Roman", 10))
        self.github_label.pack(side="bottom", padx=10, pady=10)
        self.github_label.bind("<Button-1>", lambda e: webbrowser.open_new("https://github.com/Andromeda4"))

    def browse_pdf(self):
        # Open file dialog to select a PDF file
        filepath = filedialog.askopenfilename(initialdir='/', title='Select PDF file', filetypes=(("PDF files", "*.pdf"), ("all files", "*.*")))

        # Update the filepath entry field with the selected file path
        self.filepath_entry.delete(0, tk.END)
        self.filepath_entry.insert(0, filepath)

    def summarize_pdf(self):
        # Get the file path from the entry field
        filepath = self.filepath_entry.get()

        try:
            # Open the PDF file
            pdf_file = open(filepath, 'rb')

            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfFileReader(pdf_file)

            # Initialize an empty string to store the extracted text
            text = ''

            # Loop through each page and extract the text
            for i in range(pdf_reader.getNumPages()):
                page = pdf_reader.getPage(i)
                text += page.extractText()

            # Summarize the extracted text
            # Use ChatGPT to summarize the text
            # Install transformers package and download the pre-trained model using:
            # pip install transformers
            from transformers import pipeline
            summarizer = pipeline("summarization")
            summary = summarizer(text, max_length=500, min_length=200, do_sample=False)[0]['summary_text']

            # Close the PDF file
            pdf_file.close()

            # Create a PowerPoint presentation
            prs = Presentation()

            # Add a title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Summary of PDF file"
            subtitle.text = "This PowerPoint presentation summarizes a PDF file."


            # Add a content slide with the summary text
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = "Summary"
            tf = body_shape.text_frame
            tf.text = summary

            # Save the PowerPoint presentation
            prs.save('summary.pptx')

            # Show success message box
            messagebox.showinfo("Success", "PDF summarized and PowerPoint presentation created successfully.")

        except Exception as e:
            # Show error message box
            messagebox.showerror("Error", str(e))


root = tk.Tk()
my_gui = PDFSummarizerGUI(root)
root.mainloop()
